"""
PowerPoint export skeleton for Strategy Copilot (Step 1)

- Requires: python-pptx (pip install python-pptx)
- Builds a polished 16:9 deck from your session_state
- Slides: Title, Agenda, Executive Snapshot, SWOT, Ansoff 2x2, Benchmark table, Top-5 Recommendations (Impact×Effort grid), Appendix

Usage in Streamlit (Export step):

    from export_ppt import build_ppt_from_state
    bio, fname = build_ppt_from_state(state)
    st.download_button("Download PPTX", data=bio.getvalue(), file_name=fname, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

This module is defensive: missing sections are skipped gracefully.
"""
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.util import Inches
from typing import Optional, List

def add_title(prs, title, subtitle=None):
    # helper: does a layout truly have a title placeholder?
    def has_title(layout):
        for ph in layout.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return True
        return False

    # choose a layout with a real title placeholder (fallback to first)
    layout = next((l for l in prs.slide_layouts if has_title(l)), prs.slide_layouts[0])
    slide = prs.slides.add_slide(layout)

    # set title (guard if None)
    t = slide.shapes.title
    if t is not None:
        t.text = title
        title_top = t.top
        title_height = t.height
    else:
        margin = Inches(0.6)
        title_top = Inches(0.6)
        title_height = Inches(1.0)
        tb = slide.shapes.add_textbox(margin, title_top, prs.slide_width - 2*margin, title_height)
        tb.text_frame.clear()
        tb.text_frame.paragraphs[0].add_run().text = title

    # try to use a real SUBTITLE placeholder
    placed_subtitle = False
    if subtitle:
        for ph in slide.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type == PP_PLACEHOLDER.SUBTITLE:
                ph.text = subtitle
                placed_subtitle = True
                break

        if not placed_subtitle:
            margin = Inches(0.6)
            top = title_top + title_height + Inches(0.2)
            sub_box = slide.shapes.add_textbox(margin, top, prs.slide_width - 2*margin, Inches(0.9))
            stf = sub_box.text_frame
            stf.clear()
            p2 = stf.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle

    # Accent bar (optional)
    margin = Inches(0.6)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(0.6), Inches(2.2), Inches(0.15))
    # rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_PRIMARY
    return slide

def _add_heading(slide, text: str):
    box = slide.shapes.add_textbox(MARGIN, MARGIN, W - 2*MARGIN, Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = H2_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK
    return box


def _add_bullets(slide, left, top, width, height, items: List[str]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items or []):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.level = 0
        r = p.add_run()
        r.text = str(item)
        r.font.size = BODY_SIZE
        r.font.color.rgb = COLOR_DARK
    return box


def _grid(slide, left, top, width, height):
    # Draw outer rect and cross-lines; return quadrant rects
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLOR_MED
    # Vertical and horizontal lines
    slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left + width/2, top, Inches(0.0), height).line.color.rgb = COLOR_LIGHT
    slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left, top + height/2, width, Inches(0.0)).line.color.rgb = COLOR_LIGHT
    # Return quadrants as (l, t, w, h)
    return [
        (left, top, width/2, height/2),                # Q1 TL
        (left + width/2, top, width/2, height/2),      # Q2 TR
        (left, top + height/2, width/2, height/2),     # Q3 BL
        (left + width/2, top + height/2, width/2, height/2),  # Q4 BR
    ]


def _add_small_label(slide, text: str, left, top):
    b = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.3))
    tf = b.text_frame; tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text; r.font.size = Pt(11); r.font.color.rgb = COLOR_MED

# ---------------------------- Slide builders ----------------------------

def slide_agenda(prs: Presentation, items: Optional[List[str]] = None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Agenda")
    _add_bullets(slide, MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.5), items or [
        "Inputs & Goals",
        "Framework Insights",
        "Recommendations",
        "Next Steps",
    ])
    return slide


def slide_exec_snapshot(prs: Presentation, bullets: List[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Executive Snapshot")
    _add_bullets(slide, MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.5), bullets)
    return slide


def slide_swot(prs: Presentation, swot: Dict[str, List[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "SWOT")
    box_w = (W - 3*MARGIN) / 2
    box_h = (H - 2*MARGIN - Inches(1.0)) / 2
    x1, x2 = MARGIN, MARGIN + box_w + MARGIN
    y1, y2 = Inches(1.2), Inches(1.2) + box_h + MARGIN

    def cell(title, items, x, y):
        title_box = slide.shapes.add_textbox(x, y, box_w, Inches(0.35))
        tf = title_box.text_frame; tf.clear(); p = tf.paragraphs[0]; r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = COLOR_PRIMARY
        _add_bullets(slide, x, y + Inches(0.4), box_w, box_h - Inches(0.4), items)

    cell("Strengths", swot.get("S", []), x1, y1)
    cell("Weaknesses", swot.get("W", []), x2, y1)
    cell("Opportunities", swot.get("O", []), x1, y2)
    cell("Threats", swot.get("T", []), x2, y2)
    return slide


def slide_ansoff(prs: Presentation, ansoff: Dict[str, List[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Ansoff Matrix")
    grid_left, grid_top, grid_w, grid_h = MARGIN, Inches(1.4), W - 2*MARGIN, Inches(4.6)
    quads = _grid(slide, grid_left, grid_top, grid_w, grid_h)

    labels = [
        ("Market Penetration", ansoff.get("market_penetration", [])),
        ("Product Development", ansoff.get("product_development", [])),
        ("Market Development", ansoff.get("market_development", [])),
        ("Diversification", ansoff.get("diversification", [])),
    ]
    for (title, items), (l, t, w, h) in zip(labels, quads):
        title_box = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), Inches(0.3))
        tf = title_box.text_frame; tf.clear(); p = tf.paragraphs[0]; r = p.add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = COLOR_PRIMARY
        _add_bullets(slide, l + Inches(0.1), t + Inches(0.45), w - Inches(0.2), h - Inches(0.6), items)

    _add_small_label(slide, "Existing Products → New Products", grid_left + grid_w/2 - Inches(1.2), grid_top - Inches(0.35))
    _add_small_label(slide, "Existing Markets → New Markets", grid_left - Inches(0.1), grid_top + grid_h/2 + Inches(0.05))
    return slide


def slide_benchmark(prs: Presentation, company: str, bench: Dict[str, Any]):
    table = bench.get("table") or []
    peers = bench.get("peers") or []
    if not table:
        return None
    cols = 2 + len(peers)  # Capability + company + peers
    rows = 1 + len(table)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Competitor Benchmark")

    left, top, width, height = MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.2)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table

    # Header
    hdrs = ["Capability", company] + peers
    for j, h in enumerate(hdrs):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
        cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_LIGHT

    # Body
    for i, row in enumerate(table, start=1):
        vals = [row.get("capability", "")] + [row.get(company, "")] + [row.get(p, "") for p in peers]
        for j, val in enumerate(vals):
            tbl.cell(i, j).text = str(val)

    return slide


def slide_recommendations(prs: Presentation, recs: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Top 5 Recommendations — Impact × Effort")

    grid_left, grid_top = MARGIN, Inches(1.4)
    grid_w, grid_h = W - 2*MARGIN, Inches(4.6)
    q = _grid(slide, grid_left, grid_top, grid_w, grid_h)

    # Quadrant labeling
    _add_small_label(slide, "Impact ↑", grid_left - Inches(0.05), grid_top - Inches(0.35))
    _add_small_label(slide, "Effort →", grid_left + grid_w - Inches(0.8), grid_top + grid_h + Inches(0.05))

    # Place recs into quadrants
    for idx, rec in enumerate((recs or [])[:5], start=1):
        title = rec.get("title", f"Rec {idx}")
        impact = int(rec.get("impact", 3))
        effort = int(rec.get("effort", 3))
        # Determine quadrant (simple thresholds)
        # TL (Q1): high impact (>=4), low effort (<=3)
        # TR (Q2): high impact, high effort (>3)
        # BL (Q3): low impact (<4), low effort (<=3)
        # BR (Q4): low impact, high effort (>3)
        if impact >= 4 and effort <= 3:
            quad = q[0]
        elif impact >= 4 and effort > 3:
            quad = q[1]
        elif impact < 4 and effort <= 3:
            quad = q[2]
        else:
            quad = q[3]
        l, t, w, h = quad
        box = prs.slides[-1].shapes.add_textbox(l + Inches(0.12), t + Inches(0.12), w - Inches(0.24), Inches(0.5))
        tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{idx}. {title}"; r.font.size = BODY_SIZE; r.font.color.rgb = COLOR_DARK

    # Legend
    legend = slide.shapes.add_textbox(MARGIN, grid_top + grid_h + Inches(0.2), W - 2*MARGIN, Inches(0.6))
    tfl = legend.text_frame; tfl.clear()
    p = tfl.paragraphs[0]; r = p.add_run(); r.text = "Q1: Quick Wins   Q2: Strategic Bets   Q3: Fill-ins   Q4: Long Shots"; r.font.size = Pt(12); r.font.color.rgb = COLOR_MED
    return slide


def slide_appendix_json(prs: Presentation, title: str, text: str):
    # Break long text into multiple slides
    MAX_CHARS = 2000
    chunks = [text[i:i+MAX_CHARS] for i in range(0, len(text), MAX_CHARS)] or [text]
    for i, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _add_heading(slide, f"{title}{' (cont.)' if i>1 else ''}")
        box = slide.shapes.add_textbox(MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.5))
        tf = box.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run(); r.text = chunk; r.font.name = "Courier New"; r.font.size = MONO_SIZE; r.font.color.rgb = COLOR_DARK

# ---------------------------- Orchestrator ----------------------------

def build_ppt_from_state(state: Dict[str, Any]) -> (BytesIO, str):
    """Return (pptx_bytes, filename) for download.
    Expects keys in `state`: company, product, frameworks, results, recs
    """
    company = (state.get("company") or "Company").strip()
    product = (state.get("product") or "Product").strip()
    results = state.get("results") or {}
    recs = state.get("recs") or []

    prs = Presentation()
    prs.slide_width, prs.slide_height = int(W), int(H)

    # Title
    date_str = datetime.now().strftime("%b %d, %Y")
    _add_title(prs, f"{company} × {product}", f"Strategy Snapshot — {date_str}")

    # Agenda
    slide_agenda(prs)

    # Executive Snapshot (basic heuristic based on SWOT + Ansoff presence)
    snapshot: List[str] = []
    swot = results.get("SWOT") or {}
    if any(swot.get(k) for k in ("S","W","O","T")):
        if swot.get("S"): snapshot.append(f"Strengths: {', '.join(swot['S'][:2])}")
        if swot.get("W"): snapshot.append(f"Weaknesses: {', '.join(swot['W'][:2])}")
        if swot.get("O"): snapshot.append(f"Opportunities: {', '.join(swot['O'][:2])}")
        if swot.get("T"): snapshot.append(f"Threats: {', '.join(swot['T'][:2])}")
    if results.get("Ansoff"):
        snapshot.append("Focus: Execute 1–2 high‑impact Ansoff plays next quarter.")
    if recs:
        snapshot.append(f"Top priority: {recs[0].get('title','First recommendation')}")
    slide_exec_snapshot(prs, snapshot[:6])

    # SWOT
    if swot:
        slide_swot(prs, swot)

    # Ansoff
    ansoff = results.get("Ansoff") or {}
    if ansoff:
        slide_ansoff(prs, ansoff)

    # Benchmark
    bench = results.get("Benchmark") or {}
    if bench.get("table"):
        slide_benchmark(prs, company, bench)

    # Recommendations
    if recs:
        slide_recommendations(prs, recs)

    # Appendix with raw JSON (trimmed)
    import json
    raw = json.dumps({
        "frameworks": state.get("frameworks", []),
        "results": results,
        "recs": recs,
    }, indent=2, ensure_ascii=False)
    slide_appendix_json(prs, "Appendix — Raw Analysis JSON", raw)

    # Serialize
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    safe_company = company.replace(" ", "_")
    safe_product = product.replace(" ", "_")
    fname = f"{safe_company}_{safe_product}_{datetime.now().strftime('%Y%m%d')}_strategy.pptx"
    return bio, fname

# ---------------------------- Manual test ----------------------------
if __name__ == "__main__":  # pragma: no cover
    sample_state = {
        "company": "ACME Robotics",
        "product": "Industrial IoT Sensors",
        "frameworks": ["SWOT", "Ansoff", "Benchmark"],
        "results": {
            "SWOT": {"S": ["Reliable hardware", "OEM partners"], "W": ["Brand"], "O": ["Analytics"], "T": ["Price rivals"]},
            "Ansoff": {
                "market_penetration": ["Bundle add‑ons", "Loyalty pricing"],
                "market_development": ["Enter Canada"],
                "product_development": ["Managed calibration"],
                "diversification": ["Edge‑AI module"],
            },
            "Benchmark": {
                "peers": ["Rival A", "Rival B"],
                "table": [
                    {"capability": "Sensor accuracy", "ACME Robotics": "High", "Rival A": "High", "Rival B": "Medium"},
                    {"capability": "Cloud analytics", "ACME Robotics": "Medium", "Rival A": "High", "Rival B": "Low"},
                ],
            },
        },
        "recs": [
            {"title": "OEM Bundle Program", "impact": 5, "effort": 3},
            {"title": "Managed Calibration Add‑On", "impact": 4, "effort": 3},
            {"title": "Security Proof Pack", "impact": 3, "effort": 2},
        ],
    }
    bio, name = build_ppt_from_state(sample_state)
    with open(name, "wb") as f:
        f.write(bio.getvalue())
    print("Wrote", name)
